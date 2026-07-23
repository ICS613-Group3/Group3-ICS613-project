#!/usr/bin/env python3
"""Fill Slide 2 and Slide 3 for Ivan's backend parts — clean slate."""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

SRC = "R1 Presentation.pptx"
prs = Presentation(SRC)

ACCENT = RGBColor(0x1A, 0x73, 0xE8)
DARK   = RGBColor(0x33, 0x33, 0x33)


def clear_all_paragraphs(tf):
    """Remove every <a:p> from a text frame's XML tree."""
    body_elem = tf._txBody
    for p_elem in list(body_elem.findall(qn("a:p"))):
        body_elem.remove(p_elem)


def add_body(text_frame, lines, size=13):
    """Populate a text frame with presentation-style lines."""
    clear_all_paragraphs(text_frame)
    for item in lines:
        if isinstance(item, str):
            kind, text = "body", item
        elif isinstance(item, tuple) and len(item) == 2:
            kind, text = item
        elif isinstance(item, tuple) and len(item) == 1:
            kind, text = item[0], ""
        else:
            continue

        p = text_frame.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)

        if kind == "heading":
            r.font.bold = True
            r.font.color.rgb = ACCENT
            r.font.size = Pt(15)
        elif kind == "gap":
            r.font.size = Pt(8)
        else:
            r.font.bold = False
            r.font.color.rgb = DARK


# ── Slide 2 ────────────────────────────────────────────────────────────
s2 = prs.slides[1]
# shape[0] = title, shape[1] = body text, shape[2] = ER diagram picture
title2 = s2.shapes[0].text_frame
body2  = s2.shapes[1].text_frame

clear_all_paragraphs(title2)
p = title2.add_paragraph()
r = p.add_run()
r.text = "Overview of the Domain Model and ER diagram"
r.font.size = Pt(24)
r.font.bold = True
r.font.color.rgb = DARK

add_body(body2, [
    ("heading", "8 entities, 5 core relationships"),
    "User, Tool, Reservation, Review, Invite, Photo, Notification, AdminAuditLog",
    ("gap",),
    "User owns Tools · User borrows via Reservations",
    "Reservation → Review (one per party) · Tool → Photos (1–5)",
    ("gap",),
    ("heading", "Key DB constraints"),
    "GiST EXCLUDE on (tool_id, daterange) — no double-booking",
    "UNIQUE (reservation_id, reviewer_id) — one review per party",
    "Soft-delete with active-reservation guard; hard-delete for PII",
    ("gap",),
    "6 ToolCategory · 5 ToolCondition · 6 ReservationState · 14 NotificationType",
])


# ── Slide 3 ────────────────────────────────────────────────────────────
s3 = prs.slides[2]
title3 = s3.shapes[0].text_frame
body3  = s3.shapes[1].text_frame

clear_all_paragraphs(title3)
p = title3.add_paragraph()
r = p.add_run()
r.text = "Major Accomplishments"
r.font.size = Pt(24)
r.font.bold = True
r.font.color.rgb = DARK

add_body(body3, [
    ("heading", "Backend — Ivan (BE Lead)"),
    ("gap",),
    "46 endpoints across 7 routers — all implemented and tested",
    "Auth (13), Tools (11), Reservations (10), Reviews (5), Notifications (2), Admin (4), Health (1)",
    ("gap",),
    "151 / 151 tests passing (pytest + httpx)",
    ("gap",),
    ("heading", "Infrastructure highlights"),
    "JWT access + refresh tokens with automatic rotation",
    "Rate limiting on all auth endpoints, configurable per environment",
    "GiST EXCLUDE constraint prevents overlapping reservations at DB level",
    "APScheduler: auto-cancel stale pickups, overdue escalation, token cleanup",
    "Photo uploads: multi-file, magic-byte validated, 5 MB cap, disk storage",
    ("gap",),
    ("heading", "Frontend — Yafei (FE Lead)"),
    "28 pages across auth, tool CRUD, reservations, reviews, admin, notifications",
    "AuthContext with token refresh, RequireAdmin guard, responsive layout",
    "Integration with backend in progress per API contract document",
])

prs.save(SRC)
print("OK — saved:", SRC)
