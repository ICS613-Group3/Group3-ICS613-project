from pptx import Presentation
from pptx.shapes.picture import Picture

prs = Presentation("R1 Presentation.pptx")
slide2 = prs.slides[1]
for i, sh in enumerate(slide2.shapes):
    print(f"shape[{i}]: {type(sh).__name__} name={sh.name!r}")
    if isinstance(sh, Picture):
        img = sh.image
        print(f"  content_type={img.content_type}")
        print(f"  blob_len={len(img.blob)}")
        print(f"  ext={img.ext}")
        print(f"  header: {img.blob[:20].hex()}")
        path = f"er_diagram.{img.ext}"
        with open(path, "wb") as f:
            f.write(img.blob)
        print(f"  saved to {path}")
