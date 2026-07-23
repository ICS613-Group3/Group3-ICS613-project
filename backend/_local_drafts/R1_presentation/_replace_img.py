"""Replace the ER diagram image on slide 2 with the cleaned version."""
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.util import Emu
from PIL import Image
import io

prs = Presentation("R1 Presentation.pptx")
slide2 = prs.slides[1]

# Find the picture shape
pic_shape = None
for sh in slide2.shapes:
    if isinstance(sh, Picture):
        pic_shape = sh
        break

if pic_shape is None:
    raise SystemExit("No picture found on slide 2")

# Get original position and size
left = pic_shape.left
top = pic_shape.top
width = pic_shape.width
height = pic_shape.height

print(f"Original pic: left={Emu(left)} top={Emu(top)} size={Emu(width)}x{Emu(height)}")

# Load the cleaned image, convert to PNG bytes
img = Image.open("er_diagram.png")
buf = io.BytesIO()
img.save(buf, format="PNG")
buf.seek(0)

# Replace: remove old picture, add new one at same position
pic_shape.element.getparent().remove(pic_shape.element)
slide2.shapes.add_picture(buf, left, top, width, height)

prs.save("R1 Presentation.pptx")
print("Saved R1 Presentation.pptx with cleaned ER diagram")
